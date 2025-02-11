from dotenv import load_dotenv
import google.generativeai as genai
import os
import fitz
import re
import nltk
import numpy as np
from datetime import datetime
from pptx import Presentation
from pptx.util import Pt, Inches
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from nltk.tokenize import sent_tokenize
from sklearn.feature_extraction.text import TfidfVectorizer
from sklearn.metrics.pairwise import cosine_similarity
#fi
# Load environment variables and configure
load_dotenv()
nltk.download("punkt", quiet=True)
genai.configure(api_key=os.getenv("GEMINI_API_KEY"))

class PDFProcessor:
    def __init__(self, pdf_path):
        self.pdf_path = pdf_path

    def extract_metadata(self):
        """Extract metadata from the PDF file."""
        try:
            doc = fitz.open(self.pdf_path)
            return {
                "title": doc.metadata.get("title", "Research Presentation"),
                "author": doc.metadata.get("author", "Unknown"),
                "number_of_pages": len(doc),
                "date": datetime.now().strftime("%Y-%m-%d")
            }
        except Exception as e:
            print(f"Error extracting metadata: {e}")
            return {
                "title": "Research Presentation",
                "author": "Unknown",
                "number_of_pages": 0,
                "date": datetime.now().strftime("%Y-%m-%d")
            }

    def extract_text(self):
        """Extract text content from PDF."""
        try:
            doc = fitz.open(self.pdf_path)
            text_content = {}
            for page_num in range(len(doc)):
                text_content[page_num] = doc[page_num].get_text("text")
            return text_content
        except Exception as e:
            print(f"Error extracting text: {e}")
            return {}

    def identify_sections(self, text_map):
        """Identify and extract different sections from the text."""
        sections = {}
        section_pattern = re.compile(
            r'\b(abstract|introduction|methodology|results|discussion|conclusion)\b',
            re.IGNORECASE
        )
        
        current_section = None
        section_content = ""
        
        try:
            for text in text_map.values():
                for line in text.split("\n"):
                    match = section_pattern.search(line.strip())
                    if match:
                        if current_section:
                            sections[current_section] = section_content.strip()
                        current_section = match.group(0).capitalize()
                        section_content = ""
                    else:
                        section_content += " " + line.strip()
            
            if current_section:
                sections[current_section] = section_content.strip()

            # If no sections found, create a generic one
            if not sections:
                sections["Content"] = " ".join([text for text in text_map.values()])

            return sections
        except Exception as e:
            print(f"Error identifying sections: {e}")
            return {"Content": "Error processing document sections"}

    def extract_images(self):
        """Extract images from PDF with duplicate detection."""
        try:
            doc = fitz.open(self.pdf_path)
            unique_images = set()
            images = []

            for page_num in range(len(doc)):
                page = doc[page_num]
                img_list = page.get_images(full=True)

                for img_index, img_data in enumerate(img_list):
                    try:
                        base_image = doc.extract_image(img_data[0])
                        if base_image:
                            image_bytes = base_image["image"]
                            image_hash = hash(image_bytes)
                            
                            if image_hash not in unique_images:
                                unique_images.add(image_hash)
                                img_folder = "./extracted_images"
                                os.makedirs(img_folder, exist_ok=True)
                                
                                img_path = f"{img_folder}/image_{page_num}_{img_index}.{base_image['ext']}"
                                
                                with open(img_path, "wb") as img_file:
                                    img_file.write(image_bytes)
                                
                                images.append({
                                    "page": page_num + 1,
                                    "path": img_path,
                                    "caption": f"Figure {len(images) + 1}"
                                })
                    except Exception as img_error:
                        print(f"Error processing image {img_index} on page {page_num}: {img_error}")
                        continue

            return images
        except Exception as e:
            print(f"Error extracting images: {e}")
            return []

class PresentationGenerator:
    def __init__(self, metadata, sections, images):
        self.meta = metadata
        self.sections = sections
        self.images = images
        self.prs = Presentation()
        self.theme = {
            'title_color': RGBColor(0, 32, 96),
            'text_color': RGBColor(0, 0, 0),
            'highlight_color': RGBColor(68, 114, 196),
            'background_color': RGBColor(255, 255, 255),
            'font_title': 'Arial',
            'font_body': 'Calibri',
            'title_size': Pt(44),
            'heading_size': Pt(32),
            'body_size': Pt(18)
        }

    def generate_content_with_gemini(self, section_text):
        """Generate enhanced content using Gemini AI."""
        try:
            if not os.getenv("GEMINI_API_KEY"):
                return section_text

            prompt = f"""
            Summarize the following research paper section concisely and professionally,
            maintaining key points and academic tone:
            {section_text}
            """
            
            model = genai.GenerativeModel("gemini-pro")
            response = model.generate_content(prompt)
            
            if response and response.text:
                return response.text.strip()
            return section_text
        except Exception as e:
            print(f"Error generating content with Gemini: {e}")
            return section_text

    def rank_sentences(self, text, max_sentences=5):
        """Rank and select most important sentences using TF-IDF and cosine similarity."""
        try:
            sentences = sent_tokenize(text)
            if len(sentences) <= max_sentences:
                return sentences
            
            vectorizer = TfidfVectorizer()
            tfidf_matrix = vectorizer.fit_transform(sentences)
            similarity_matrix = cosine_similarity(tfidf_matrix, tfidf_matrix)
            sentence_ranks = similarity_matrix.sum(axis=1)
            ranked_sentence_indexes = np.argsort(sentence_ranks)[::-1]
            
            return [sentences[i] for i in ranked_sentence_indexes[:max_sentences]]
        except Exception as e:
            print(f"Error ranking sentences: {e}")
            return text.split('. ')[:max_sentences]

    def apply_text_formatting(self, shape, is_title=False):
        """Apply consistent text formatting to shapes."""
        paragraph = shape.text_frame.paragraphs[0]
        paragraph.alignment = PP_ALIGN.CENTER if is_title else PP_ALIGN.LEFT
        
        font = paragraph.font
        font.name = self.theme['font_title'] if is_title else self.theme['font_body']
        font.size = self.theme['title_size'] if is_title else self.theme['body_size']
        font.color.rgb = self.theme['title_color'] if is_title else self.theme['text_color']

    def create_title_slide(self):
        """Create the title slide."""
        slide_layout = self.prs.slide_layouts[0]
        slide = self.prs.slides.add_slide(slide_layout)
        
        title = slide.shapes.title
        subtitle = slide.placeholders[1]
        
        title.text = self.meta.get("title", "Research Presentation")
        subtitle.text = f"Author: {self.meta['author']}\nDate: {self.meta['date']}"
        
        self.apply_text_formatting(title, is_title=True)
        self.apply_text_formatting(subtitle)

    def create_section_slide(self, title, content):
        """Create a slide for each section."""
        slide_layout = self.prs.slide_layouts[1]
        slide = self.prs.slides.add_slide(slide_layout)
        
        title_shape = slide.shapes.title
        content_shape = slide.placeholders[1]
        
        title_shape.text = title
        self.apply_text_formatting(title_shape, is_title=True)
        
        enhanced_content = self.generate_content_with_gemini(content)
        ranked_sentences = self.rank_sentences(enhanced_content)
        content_shape.text = "\n• " + "\n• ".join(ranked_sentences)
        self.apply_text_formatting(content_shape)

    def create_image_slide(self, image_info):
        """Create a slide for each image."""
        slide_layout = self.prs.slide_layouts[1]
        slide = self.prs.slides.add_slide(slide_layout)
        
        title_shape = slide.shapes.title
        title_shape.text = image_info['caption']
        self.apply_text_formatting(title_shape, is_title=True)
        
        # Calculate image position and size
        left = Inches(1)
        top = Inches(1.5)
        width = Inches(8)
        height = Inches(5)
        
        try:
            slide.shapes.add_picture(
                image_info['path'],
                left, top, width, height
            )
        except Exception as e:
            print(f"Error adding image to slide: {e}")

    def create_presentation(self):
        """Create the complete presentation."""
        try:
            # Create title slide
            self.create_title_slide()
            
            # Create content slides
            for section, content in self.sections.items():
                self.create_section_slide(section, content)
            
            # Create image slides
            for image in self.images:
                self.create_image_slide(image)
                
        except Exception as e:
            print(f"Error creating presentation: {e}")

    def save(self, output_path="./output/presentation.pptx"):
        """Save the presentation to file."""
        try:
            os.makedirs(os.path.dirname(output_path), exist_ok=True)
            self.prs.save(output_path)
            print(f"Presentation saved successfully at {output_path}")
            return output_path
        except Exception as e:
            print(f"Error saving presentation: {e}")
            raise

def process_pdf_to_ppt(pdf_path, output_path=None):
    """Main function to process PDF and generate PPT."""
    try:
        # Initialize PDF processor
        processor = PDFProcessor(pdf_path)
        
        # Extract content
        metadata = processor.extract_metadata()
        text_map = processor.extract_text()
        sections = processor.identify_sections(text_map)
        images = processor.extract_images()
        
        # Generate presentation
        presentation = PresentationGenerator(metadata, sections, images)
        presentation.create_presentation()
        
        # Save presentation
        if output_path is None:
            timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
            output_path = f"./output/presentation_{timestamp}.pptx"
        
        return presentation.save(output_path)
    
    except Exception as e:
        print(f"Error processing PDF to PPT: {e}")